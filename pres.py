from pptx import Presentation
import pandas as pd
import json


file = './test.pptx'
config_file = './conf.json'
outfile = './newin.json'

general_data = json.load(open(config_file, 'r'))
file = f'./{general_data["name"]}'
ppt=Presentation(file)

notes = []
titles = []
page_info ={}


for page, slide in enumerate(ppt.slides):
    # this is the notes that doesn't appear on the ppt slide,
    # but really the 'presenter' note. 
    textNote = slide.notes_slide.notes_text_frame.text
    try:
        title = slide.shapes.title.text
    except:
        title = f"Folie {str(page+1)}"
    page_info[f'p{page}'] = {
        "titel":title,
        "note":textNote,
        }
        
for i in page_info.keys():
    page_info[i]["time"] = general_data["page_time"] if not general_data["page_time"] else general_data["tot_time"]/ len(page_info.keys())

# generate settings for watch
## general
settings = {
    "name": general_data["name"],
    "pages": len(page_info.keys()),
    "nat_scroll": general_data["natural_scrolling_direction"],
    "right_arm":general_data["right_arm"],
    "page_time" : general_data["page_time"],
    "total_time":general_data["tot_time"],
    "final_quote": general_data["final_statement"]
}

print(settings)
print(page_info)
with open(outfile, 'w', encoding='utf-8') as file:
    json.dump({"settings":settings, "page_info":page_info}, file)
