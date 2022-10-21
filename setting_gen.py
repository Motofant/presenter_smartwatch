from pptx import Presentation
import PyPDF2 as pypdf
#import odf
import os
import json

file  = "./Test.pptx"

#file = "./Verteidigung_presentation.pdf"
"""
file = "./test_new.pptm"
file = "./test.odp"
"""
settings_file = "./newin2.json"
full_support = [".pptx",".pptm"]
part_support = [".pdf",".odg"]
total_dictionary = {}

# get file extension
extension = os.path.splitext(file)[-1]
print(extension)

if extension in full_support:
    print("full support")
    ppt = Presentation(file)
    total_dictionary = {"settings":{
        "name":file,
        "page":len(ppt.slides),
        "gesture" : "false",
        "comments" : "true",
        "page_time" : "true"
    },"page_info":{

    }}

    id = 0
    for page, slide in enumerate(ppt.slides):
        page_id = f"p{id}"

        try:
            title = slide.shapes.title.text
        except:
            title = f"Folie {str(page+1)}"

        note = slide.notes_slide.notes_text_frame.text

        total_dictionary["page_info"][page_id] = {
            "titel" : title,
            "note" : note,
            "time" : "false"   
        }


        id+= 1 

elif extension == ".pdf":
    pdf=pypdf.PdfFileReader(open(file,'rb'))
    num_pages = pdf.numPages

    total_dictionary = {"settings":{
        "name":file,
        "page": num_pages,
        "gesture" : "true",
        "comments" : "false",
        "page_time" : "true"
    },"page_info":{

    }}

    for i in range(0,num_pages-1):
        page_id = f"p{i}"

        title = f"Folie {str(i+1)}"

        total_dictionary["page_info"][page_id] = {
            "titel" : title,
            "note" : "",
            "time" : "false"   
        }
    """
elif extension == ".odp":
    odf.opendocument.load(file)
    """
else:
    print("file format not supported")
    print("yet?")
with open (settings_file,"w") as outfile:
    json.dump(total_dictionary,outfile)